import configparser
from pptx.util import Pt
from pptx.util import Inches
import os
import tkinter as tk
from tkinter import filedialog, messagebox, Listbox, MULTIPLE, Toplevel, Canvas, Scrollbar, Menu, StringVar
from tkinter.ttk import Progressbar, Label, Entry, Button, Checkbutton
from PIL import Image, ImageTk, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
import subprocess
from pathlib import Path
from io import BytesIO
import threading

class ConfigManager:
    def __init__(self, config_path='config.ini'):
        self.config_path = config_path
        self.config = configparser.ConfigParser()
        self.load_config()

    def load_config(self):
        if not os.path.exists(self.config_path):
            self.create_default_config()
        self.config.read(self.config_path)

        required_options = {
            'order_of_service': 'Opening Prayer, Song, Song, Giving, Song, Lords Supper, Song, Scripture Reading, Sermon, Song, Announcements, Song, Closing Prayer',
            'default_song_search_path': '/home/church/Music/PaperlessHymnal(Vol.1-13)/4x3',
            'add_song_title_slides': 'true',
            'announcement_slides_path': '/home/church/Desktop/Announcement slides.pptx',
            'temp_directory': '/tmp/ppt_combiner_temp',
            'song_title_slide_title': 'Let us sing:\n',
            'title_slide_trim_chars': '-,_,.ppt',
            'default_directory': '/home/church/Music/PaperlessHymnal(Vol.1-13)/4x3/',
            'unoconv_path': '/usr/bin/unoconv'
        }

        if 'Settings' not in self.config:
            self.config['Settings'] = {}

        updated = False
        for option, default_value in required_options.items():
            if option not in self.config['Settings']:
                self.config['Settings'][option] = default_value
                updated = True

        if updated:
            with open(self.config_path, 'w') as configfile:
                self.config.write(configfile)

    def create_default_config(self):
        self.config['Settings'] = {
            'order_of_service': 'Opening Prayer, Song, Song, Giving, Song, Lords Supper, Song, Scripture Reading, Sermon, Song, Announcements, Song, Closing Prayer',
            'default_song_search_path': '/home/church/Music/PaperlessHymnal(Vol.1-13)/4x3',
            'add_song_title_slides': 'true',
            'announcement_slides_path': '/home/church/Desktop/Announcement slides.pptx',
            'temp_directory': '/tmp/ppt_combiner_temp',
            'song_title_slide_title': 'Let us sing:\n',
            'title_slide_trim_chars': '-,_,.ppt',
            'default_directory': '/home/church/Music/PaperlessHymnal(Vol.1-13)/4x3/',
            'title_slide_directory': '/path/to/title/slides'
        }
        with open(self.config_path, 'w') as configfile:
            self.config.write(configfile)

    def get(self, key):
        return self.config['Settings'].get(key)

    def set(self, key, value):
        self.config['Settings'][key] = value
        with open(self.config_path, 'w') as configfile:
            self.config.write(configfile)

class OptionsPage:
    def __init__(self, root, config_manager):
        self.root = root
        self.config_manager = config_manager
        self.options_dialog = Toplevel(root)
        self.options_dialog.title("Settings")

        # Configuration options
        Label(self.options_dialog, text="Order of Service:").grid(row=0, column=0, padx=10, pady=5, sticky='w')
        self.order_of_service_entry = Entry(self.options_dialog, width=50)
        self.order_of_service_entry.insert(0, self.config_manager.get('order_of_service'))
        self.order_of_service_entry.grid(row=0, column=1, padx=10, pady=5)

        Label(self.options_dialog, text="Default Song Search Path:").grid(row=1, column=0, padx=10, pady=5, sticky='w')
        self.default_song_search_path_entry = Entry(self.options_dialog, width=50)
        self.default_song_search_path_entry.insert(0, self.config_manager.get('default_song_search_path'))
        self.default_song_search_path_entry.grid(row=1, column=1, padx=10, pady=5)

        Label(self.options_dialog, text="Announcement Slides Path:").grid(row=2, column=0, padx=10, pady=5, sticky='w')
        self.announcement_slides_path_entry = Entry(self.options_dialog, width=50)
        self.announcement_slides_path_entry.insert(0, self.config_manager.get('announcement_slides_path'))
        self.announcement_slides_path_entry.grid(row=2, column=1, padx=10, pady=5)

        Label(self.options_dialog, text="Temporary Directory:").grid(row=3, column=0, padx=10, pady=5, sticky='w')
        self.temp_directory_entry = Entry(self.options_dialog, width=50)
        self.temp_directory_entry.insert(0, self.config_manager.get('temp_directory'))
        self.temp_directory_entry.grid(row=3, column=1, padx=10, pady=5)

        Label(self.options_dialog, text="Song Title Slide Title:").grid(row=4, column=0, padx=10, pady=5, sticky='w')
        self.song_title_slide_title_entry = Entry(self.options_dialog, width=50)
        self.song_title_slide_title_entry.insert(0, self.config_manager.get('song_title_slide_title'))
        self.song_title_slide_title_entry.grid(row=4, column=1, padx=10, pady=5)

        Label(self.options_dialog, text="Title Slide Trim Characters:").grid(row=5, column=0, padx=10, pady=5, sticky='w')
        self.title_slide_trim_chars_entry = Entry(self.options_dialog, width=50)
        self.title_slide_trim_chars_entry.insert(0, self.config_manager.get('title_slide_trim_chars'))
        self.title_slide_trim_chars_entry.grid(row=5, column=1, padx=10, pady=5)

        self.add_song_title_slides_var = StringVar(value=self.config_manager.get('add_song_title_slides'))

        Label(self.options_dialog, text="Add Song Title Slides:").grid(row=6, column=0, padx=10, pady=5, sticky='w')
        Checkbutton(self.options_dialog, text="Enable", variable=self.add_song_title_slides_var, onvalue='true', offvalue='false').grid(row=6, column=1, padx=10, pady=5, sticky='w')

        Label(self.options_dialog, text="Title Slide Directory:").grid(row=7, column=0, padx=10, pady=5, sticky='w')
        self.title_slide_directory_entry = Entry(self.options_dialog, width=50)
        title_slide_directory = self.config_manager.get('title_slide_directory') or ""
        self.title_slide_directory_entry.insert(0, title_slide_directory)
        self.title_slide_directory_entry.grid(row=7, column=1, padx=10, pady=5)

        Button(self.options_dialog, text="Save", command=self.save_settings).grid(row=8, column=0, columnspan=2, pady=10)

    def save_settings(self):
        self.config_manager.set('order_of_service', self.order_of_service_entry.get())
        self.config_manager.set('default_song_search_path', self.default_song_search_path_entry.get())
        self.config_manager.set('announcement_slides_path', self.announcement_slides_path_entry.get())
        self.config_manager.set('temp_directory', self.temp_directory_entry.get())
        self.config_manager.set('song_title_slide_title', self.song_title_slide_title_entry.get())
        self.config_manager.set('title_slide_trim_chars', self.title_slide_trim_chars_entry.get())
        self.config_manager.set('add_song_title_slides', self.add_song_title_slides_var.get())
        self.config_manager.set('title_slide_directory', self.title_slide_directory_entry.get())
        self.options_dialog.destroy()

class PPTCombinerApp:
    def __init__(self, root, config_manager):
        self.root = root
        self.root.title("Church Service Organizer")
        
        self.config_manager = config_manager
        self.order_of_service = config_manager.get('order_of_service').split(', ')
        self.default_directory = config_manager.get('default_song_search_path')
        self.unoconv_path = config_manager.get('unoconv_path')
        self.add_song_title_slides = config_manager.get('add_song_title_slides') == 'true'
        self.song_title_slide_title = config_manager.get('song_title_slide_title')
        self.title_slide_trim_chars = config_manager.get('title_slide_trim_chars').split(',')
        self.all_ppt_files = []
        self.filtered_ppt_files = []
        self.playlist = []

        self.setup_ui()
        self.load_order_of_service()  # Load the order of worship elements into the playlist manager window
        self.search_files()  # Run the search in the default directory on boot
    
    def setup_ui(self):
        # Create menu
        menu = Menu(self.root)
        self.root.config(menu=menu)
        
        options_menu = Menu(menu, tearoff=0)
        menu.add_cascade(label="Options", menu=options_menu)
        options_menu.add_command(label="Settings", command=self.open_options_dialog)

        frame = tk.Frame(self.root)
        frame.pack(fill=tk.BOTH, expand=True)

        left_frame = tk.Frame(frame)
        left_frame.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

        right_frame = tk.Frame(frame)
        right_frame.pack(side=tk.RIGHT, fill=tk.BOTH, expand=True)

        # Directory selection
        self.dir_label = tk.Label(left_frame, text="Select Directory:", anchor="w")
        self.dir_label.pack(fill="x")
        
        self.dir_entry = tk.Entry(left_frame, width=50)
        self.dir_entry.insert(0, self.default_directory)
        self.dir_entry.pack(fill="x")
        
        self.browse_button = tk.Button(left_frame, text="Browse", command=self.browse_directory)
        self.browse_button.pack(fill="x")
        
        # Search input
        self.search_label = tk.Label(left_frame, text="Search:", anchor="w")
        self.search_label.pack(fill="x")
        
        self.search_entry = tk.Entry(left_frame, width=50)
        self.search_entry.pack(fill="x")
        self.search_entry.bind('<KeyRelease>', self.filter_files)
        
        # Add to playlist button
        self.add_to_playlist_button = tk.Button(left_frame, text="Add to Playlist", command=self.add_to_playlist)
        self.add_to_playlist_button.pack(fill="x")
        
        # File list
        #self.file_listbox = Listbox(left_frame, selectmode=tk.MULTIPLE, width=100, height=20)
        self.shift_click_start_index = None  # Track the first click for shift selection

        self.file_listbox = Listbox(left_frame, selectmode=tk.BROWSE, width=100, height=20)
        self.file_listbox.pack(fill="both", expand=True)
        self.file_listbox.bind('<Double-1>', self.add_to_playlist)
        self.file_listbox.bind('<Control-1>', self.preview_slide)

        # Update the file listbox to support multi-selection mode
        # Allow multi-selection of PowerPoints for adding multiple files to the playlist
        #self.add_to_playlist_button = tk.Button(left_frame, text="Add to Playlist", command=self.add_to_playlist)

        
        # Playlist
        self.playlist_label = tk.Label(right_frame, text="Playlist:", anchor="w")
        self.playlist_label.pack(fill="x")
        
        self.playlist_listbox = Listbox(right_frame, selectmode=tk.BROWSE, width=100, height=20)
        self.playlist_listbox.pack(fill="both", expand=True)

        self.move_up_button = tk.Button(right_frame, text="Move Up", command=self.move_up)
        self.move_up_button.pack(fill="x")

        self.move_down_button = tk.Button(right_frame, text="Move Down", command=self.move_down)
        self.move_down_button.pack(fill="x")

        self.remove_button = tk.Button(right_frame, text="Remove", command=self.remove_from_playlist)
        self.remove_button.pack(fill="x")
        
        # Save and Convert button
        self.save_convert_button = tk.Button(right_frame, text="Save and Convert", command=self.save_and_convert)
        self.save_convert_button.pack(fill="x")

    def on_shift_click_select(self, event):
        """Handles shift-click selection to select multiple files in BROWSE mode."""
        current_index = self.file_listbox.nearest(event.y)
        if event.state & 0x0001:  # Check if Shift key is held
            # Set the range if shift-clicking
            if self.shift_click_start_index is not None:
                # Select all files in the range from the first shift-click to the current click
                start = min(self.shift_click_start_index, current_index)
                end = max(self.shift_click_start_index, current_index)
                self.file_listbox.selection_clear(0, tk.END)
                for i in range(start, end + 1):
                    self.file_listbox.selection_set(i)
            else:
                self.shift_click_start_index = current_index
                self.file_listbox.selection_set(current_index)
        else:
            # Set initial click position if no Shift key is held
            self.shift_click_start_index = current_index



    def load_order_of_service(self):
        for item in self.order_of_service:
            if item.lower() == "song":
                self.playlist.append(item)
                self.playlist_listbox.insert(tk.END, item)
            else:
                self.playlist.append(item)
                self.playlist_listbox.insert(tk.END, item)

    def browse_directory(self):
        directory = filedialog.askdirectory(initialdir=self.default_directory)
        if directory:
            self.dir_entry.delete(0, tk.END)
            self.dir_entry.insert(0, directory)
            self.search_files()  # Update the search results when a new directory is selected
    
    def search_files(self):
        directory = self.dir_entry.get()
        if not os.path.isdir(directory):
            messagebox.showerror("Error", "Invalid directory")
            return
        
        self.all_ppt_files = self.recursive_search(directory)
        self.filtered_ppt_files = self.all_ppt_files
        self.update_file_listbox(self.filtered_ppt_files)
    
    def recursive_search(self, directory):
        matches = []
        for root, _, files in os.walk(directory):
            for file in files:
                if file.endswith(('.ppt', '.pptx')):
                    matches.append(os.path.join(root, file))
        return sorted(matches)  # Sort the results
    
    def filter_files(self, event):
        query = self.search_entry.get().lower()
        self.filtered_ppt_files = [file for file in self.all_ppt_files if query in os.path.basename(file).lower()]
        self.update_file_listbox(self.filtered_ppt_files)
    
    def update_file_listbox(self, files):
        self.file_listbox.delete(0, tk.END)
        for file in files:
            self.file_listbox.insert(tk.END, os.path.basename(file))


    def add_to_playlist(self, event=None):
        """Add all matching song parts (title slide, verses, and chorus) to the playlist in place of the first 'Song' entry."""
        selected_index = self.file_listbox.curselection()
        if not selected_index:
            return
    
        # Get the selected file and determine its base song name
        selected_file = self.filtered_ppt_files[selected_index[0]]
        song_name = self.extract_song_name(selected_file)
    
        # Find all matching song parts in the directory
        matching_files = [
            file for file in self.all_ppt_files
            if self.extract_song_name(file) == song_name
        ]
    
        # Separate verses and chorus, and sort verses by their order in the filename
        verses = sorted([f for f in matching_files if "verse" in f.lower()])
        chorus = next((f for f in matching_files if "chorus" in f.lower()), None)

        # Locate the first "Song" entry to replace in the playlist
        song_index = next((i for i, item in enumerate(self.playlist) if item.lower() == "song"), None)
        insert_position = song_index if song_index is not None else len(self.playlist)
    
        # Clear the current "Song" placeholder if found
        if song_index is not None:
            self.playlist.pop(song_index)
  

        # Add the hymn slides, inserting chorus after each verse
        for verse in verses:
            self.playlist.insert(insert_position, verse)
            self.playlist_listbox.insert(insert_position, os.path.basename(verse))
            insert_position += 1
        
            if chorus:
                self.playlist.insert(insert_position, chorus)
                self.playlist_listbox.insert(insert_position, os.path.basename(chorus))
                insert_position += 1

        self.refresh_playlist_display()



    def find_title_slide(self, song_name):
        """Search recursively for a title slide PNG file that matches the song name with '- Title' suffix."""
        title_slide_dir = self.config_manager.get('title_slide_directory')
        if not title_slide_dir or not os.path.isdir(title_slide_dir):
            return None

        # Expected title slide name
        expected_title = f"{song_name} - Title.png"
    
        # Use os.walk to search recursively
        for root, _, files in os.walk(title_slide_dir):
            for file in files:
                if file.lower() == expected_title.lower():
                    return os.path.join(root, file)

        # If no matching title slide is found, return None
        return None


    

    def extract_song_name(self, file):
        """Extract base name of the song from file name, ignoring verse/chorus distinctions."""
        base_name = os.path.basename(file).split('.')[0]
        # Removing verse or chorus identifiers, e.g., "SongName - Verse 1" to "SongName"
        return base_name.rsplit('-', 1)[0].strip()

    
    def move_up(self):
        """Move the selected hymn group (verses and chorus) up in the playlist."""
        selected_index = self.playlist_listbox.curselection()
        if not selected_index:
            return
    
        # Determine the base song name for the selected item
        selected_file = self.playlist[selected_index[0]]
        song_name = self.extract_song_name(selected_file)
    
        # Find all indices of items that match the song name
        group_indices = [i for i, item in enumerate(self.playlist) if self.extract_song_name(item) == song_name]
    
        # Only proceed if the group can be moved up
        if group_indices[0] > 0:
            # Move each item in the group up by one position
            for i in group_indices:
                self.playlist[i], self.playlist[i-1] = self.playlist[i-1], self.playlist[i]
        
            # Update the playlist display
            self.refresh_playlist_display()
            self.playlist_listbox.select_set(group_indices[0] - 1)

    def move_down(self):
        """Move the selected hymn group (verses and chorus) down in the playlist."""
        selected_index = self.playlist_listbox.curselection()
        if not selected_index:
            return
    
        # Determine the base song name for the selected item
        selected_file = self.playlist[selected_index[0]]
        song_name = self.extract_song_name(selected_file)
    
        # Find all indices of items that match the song name
        group_indices = [i for i, item in enumerate(self.playlist) if self.extract_song_name(item) == song_name]
    
        # Only proceed if the group can be moved down
        if group_indices[-1] < len(self.playlist) - 1:
            # Move each item in the group down by one position
            for i in reversed(group_indices):
                self.playlist[i], self.playlist[i+1] = self.playlist[i+1], self.playlist[i]
        
            # Update the playlist display
            self.refresh_playlist_display()
            self.playlist_listbox.select_set(group_indices[0] + 1)

    def refresh_playlist_display(self):
        """Refreshes the playlist display in the Listbox."""
        self.playlist_listbox.delete(0, tk.END)
        for item in self.playlist:
            self.playlist_listbox.insert(tk.END, os.path.basename(item))

    
    def remove_from_playlist(self):
        """Remove the selected hymn group (verses and chorus) from the playlist."""
        selected_index = self.playlist_listbox.curselection()
        if not selected_index:
            return

        # Determine the base song name for the selected item
        selected_file = self.playlist[selected_index[0]]
        song_name = self.extract_song_name(selected_file)
    
        # Find all indices of items that match the song name
        group_indices = [i for i, item in enumerate(self.playlist) if self.extract_song_name(item) == song_name]
    
        # Remove items in the group from the playlist in reverse order to avoid index shifting
        for i in reversed(group_indices):
            del self.playlist[i]
    
        # Refresh the playlist display
        self.refresh_playlist_display()
    
    def save_and_convert(self):
        desktop_path = str(Path.home() / "Desktop")
        save_path = filedialog.asksaveasfilename(initialdir=desktop_path, defaultextension=".pptx", filetypes=[("PowerPoint files", "*.pptx")])
        if not save_path:
            return

        # Create a new window for the loading indicator
        loading_window = Toplevel(self.root)
        loading_window.title("Saving and Converting")

        progress = Progressbar(loading_window, mode='indeterminate')
        progress.pack(pady=20)
        progress.start()

        def perform_conversion():
            self.combine_presentations(save_path)
            progress.stop()
            loading_window.destroy()

        threading.Thread(target=perform_conversion).start()
    
    def convert_ppt_to_pptx(self, file_path):
        if file_path.endswith('.pptx'):
            return file_path  # No conversion needed
        pptx_path = file_path.replace('.ppt', '.pptx')
        try:
            subprocess.run([self.unoconv_path, '-f', 'pptx', file_path], check=True)
        except FileNotFoundError:
            messagebox.showerror("Error", "unoconv not found. Please install unoconv to convert .ppt files.")
            return None
        except subprocess.CalledProcessError as e:
            messagebox.showerror("Error", f"Conversion failed: {e}")
            return None
        return pptx_path



    def combine_presentations(self, save_path):
        """Combine presentations, adding a title slide image if available, and save to save_path."""
        combined_ppt = Presentation()
        temp_directory = self.config_manager.get('temp_directory')

        # Ensure temporary directory exists
        if not os.path.exists(temp_directory):
            os.makedirs(temp_directory)

        # Set slide dimensions
        combined_ppt.slide_width = Inches(13.33)
        combined_ppt.slide_height = Inches(7.5)

        # Loop through playlist items
        title_added = False  # Track if title slide has been added for the current hymn
        current_hymn = ""
        for item in self.playlist:
            # If this is a hymn, add the title slide image only once before the first verse or chorus
            if os.path.exists(item) and item.endswith(('.ppt', '.pptx')):
                song_name = self.extract_song_name(item)
                if current_hymn != song_name:
                    title_added = False
                current_hymn = song_name 
                if not title_added:
                    title_slide_path = self.find_title_slide(song_name)     
                    # Add title slide if it exists and hasn't been added yet
                    if title_slide_path:
                        self.add_image_slide(combined_ppt, title_slide_path)
                        title_added = True  # Mark title as added for this hymn
    
                # Add the hymn's slides to the presentation
                ppt = Presentation(item)
                for slide in ppt.slides:
                    self.add_slide_to_presentation(combined_ppt, slide)
            else:
    
                # Add non-hymn slide with centered text and no placeholder
                slide = combined_ppt.slides.add_slide(combined_ppt.slide_layouts[5])

                # Remove any placeholder elements ("Click to add title" boxes)
                for shape in slide.shapes:
                    if shape.is_placeholder:
                        sp = slide.shapes._spTree.remove(shape._element)

                # Add a text box with centered text
                text_box = slide.shapes.add_textbox(
                    Inches(1), Inches(1.5),  # Adjust positioning
                    combined_ppt.slide_width - Inches(2), combined_ppt.slide_height - Inches(3)
                )
                text_frame = text_box.text_frame
                text_frame.text = item
                text_frame.paragraphs[0].font.size = Pt(80)
                text_frame.paragraphs[0].font.bold = True
                text_box.text_anchor = "middle"  # Center text within text box
                text_box.left = int((combined_ppt.slide_width - text_box.width) / 2)
                text_box.top = int((combined_ppt.slide_height - text_box.height) / 2)  # Center vertically on slide

        combined_ppt.save(save_path)
        messagebox.showinfo("Success", f"Combined presentation saved to {save_path}")


    def add_image_slide(self, presentation, image_path):
        """Add an image as a full-slide background in the given presentation."""
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        img = Image.open(image_path)
        image_stream = BytesIO()
        img.save(image_stream, format='PNG')

        # Add image to slide
        slide.shapes.add_picture(image_stream, 0, 0, width=presentation.slide_width, height=presentation.slide_height)


    def add_slide_to_presentation(self, combined_ppt, slide):
        slide_copy = combined_ppt.slides.add_slide(combined_ppt.slide_layouts[5])

        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                image_stream = shape.image.blob
                image_file = BytesIO(image_stream)
                image = Image.open(image_file)
                temp_image_path = os.path.join(self.config_manager.get('temp_directory'), "temp_image.png")
                image.save(temp_image_path)
                slide_copy.shapes.add_picture(temp_image_path, shape.left, shape.top, shape.width, shape.height)

    def preview_slide(self, event):
        selected_index = self.file_listbox.nearest(event.y)
        selected_file = self.filtered_ppt_files[selected_index]
        file_name = os.path.basename(selected_file)

        preview_window = Toplevel(self.root)
        preview_window.title(file_name)
        
        progress = Progressbar(preview_window, mode='indeterminate')
        progress.pack(pady=20)
        progress.start()

        def load_preview():
            pptx_file = self.convert_ppt_to_pptx(selected_file)
            if not pptx_file:
                progress.stop()
                progress.pack_forget()
                return

            ppt = Presentation(pptx_file)
            slide = ppt.slides[0]  # Preview the first slide

            img = self.get_slide_image(slide)
            if img:
                img_tk = ImageTk.PhotoImage(img)
                canvas = Canvas(preview_window, width=img.width, height=img.height)
                canvas.create_image(0, 0, anchor='nw', image=img_tk)
                canvas.img = img_tk
                h_scroll = Scrollbar(preview_window, orient='horizontal', command=canvas.xview)
                v_scroll = Scrollbar(preview_window, orient='vertical', command=canvas.yview)
                canvas.config(scrollregion=(0, 0, img.width, img.height), xscrollcommand=h_scroll.set, yscrollcommand=v_scroll.set)
                
                canvas.pack(side='left', fill='both', expand=True)
                h_scroll.pack(side='bottom', fill='x')
                v_scroll.pack(side='right', fill='y')

                canvas.bind("<MouseWheel>", lambda event: self.on_mouse_wheel(event, canvas))
                canvas.bind("<Shift-MouseWheel>", lambda event: self.on_shift_mouse_wheel(event, canvas))
            else:
                messagebox.showerror("Error", "No image found in the first slide.")
            
            progress.stop()
            progress.pack_forget()

        threading.Thread(target=load_preview).start()

    def on_mouse_wheel(self, event, canvas):
        if event.delta:
            canvas.yview_scroll(int(-1 * (event.delta / 120)), "units")
        elif event.num == 5:
            canvas.yview_scroll(1, "units")
        elif event.num == 4:
            canvas.yview_scroll(-1, "units")

    def on_shift_mouse_wheel(self, event, canvas):
        if event.delta:
            canvas.xview_scroll(int(-1 * (event.delta / 120)), "units")
        elif event.num == 5:
            canvas.xview_scroll(1, "units")
        elif event.num == 4:
            canvas.xview_scroll(-1, "units")

    def get_slide_image(self, slide, max_size=(1200, 900)):
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                image_stream = shape.image.blob
                image_file = BytesIO(image_stream)
                image = Image.open(image_file)
                
                image.thumbnail(max_size, Image.LANCZOS)  # Scale down the image if it's too large
                
                return image
        return None

    def open_options_dialog(self):
        OptionsPage(self.root, self.config_manager)

if __name__ == "__main__":
    config_manager = ConfigManager()
    root = tk.Tk()
    app = PPTCombinerApp(root, config_manager)
    root.mainloop()

